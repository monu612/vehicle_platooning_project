import os
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ASSETS_DIR = "ppt_assets"
os.makedirs(ASSETS_DIR, exist_ok=True)

# THEME COLORS
BG_COLOR = "#0a0f1a"
TEXT_COLOR = "white"
CYAN = "#00ffff"
PINK = "#ff007f"

def setup_plot(figsize=(8, 4)):
    fig, ax = plt.subplots(figsize=figsize, facecolor=BG_COLOR)
    ax.set_facecolor(BG_COLOR)
    ax.axis('off')
    return fig, ax

def save_fig(fig, filename):
    filepath = os.path.join(ASSETS_DIR, filename)
    fig.savefig(filepath, facecolor=BG_COLOR, bbox_inches='tight', dpi=300)
    plt.close(fig)
    return filepath

# --- GENERATE IMAGES ---

# 1. Platooning Diagram
fig, ax = setup_plot(figsize=(10, 2))
for i in range(7):
    color = CYAN if i == 0 else "white"
    rect = patches.Rectangle((i*2, 0), 1.5, 0.8, linewidth=2, edgecolor=color, facecolor='none')
    ax.add_patch(rect)
    label = "Leader (M)" if i == 0 else f"S{i}"
    ax.text(i*2 + 0.75, 0.4, label, color=color, ha='center', va='center', fontsize=12, fontweight='bold')
    if i < 6:
        # Connection lines
        ax.plot([i*2 + 1.5, i*2 + 2], [0.4, 0.4], color=CYAN, linestyle='--', linewidth=2)
ax.set_xlim(-0.5, 13.5)
ax.set_ylim(-1, 2)
platoon_img = save_fig(fig, "platoon.png")

# 2. Communication Failure Diagram
fig, ax = setup_plot(figsize=(10, 2))
for i in range(7):
    color = CYAN if i == 0 else "white"
    rect = patches.Rectangle((i*2, 0), 1.5, 0.8, linewidth=2, edgecolor=color, facecolor='none')
    ax.add_patch(rect)
    label = "Leader (M)" if i == 0 else f"S{i}"
    ax.text(i*2 + 0.75, 0.4, label, color=color, ha='center', va='center', fontsize=12, fontweight='bold')
    if i < 6:
        if i == 2:
            ax.plot([i*2 + 1.5, i*2 + 2], [0.4, 0.4], color=PINK, linestyle='--', linewidth=2)
            ax.text(i*2 + 1.75, 0.4, "X", color=PINK, ha='center', va='center', fontsize=20, fontweight='bold')
        else:
            ax.plot([i*2 + 1.5, i*2 + 2], [0.4, 0.4], color=CYAN, linestyle='--', linewidth=2)
ax.set_xlim(-0.5, 13.5)
ax.set_ylim(-1, 2)
failure_img = save_fig(fig, "failure.png")

# 3. Star Topology
fig, ax = setup_plot(figsize=(6, 6))
G_star = nx.star_graph(6)
pos_star = nx.spring_layout(G_star, seed=42)
nx.draw_networkx_nodes(G_star, pos_star, node_color=BG_COLOR, edgecolors=CYAN, node_size=1000, linewidths=2, ax=ax)
nx.draw_networkx_edges(G_star, pos_star, edge_color=CYAN, width=2, ax=ax)
labels = {0: "M", 1: "S1", 2: "S2", 3: "S3", 4: "S4", 5: "S5", 6: "S6"}
nx.draw_networkx_labels(G_star, pos_star, labels, font_color="white", font_weight="bold", ax=ax)
ax.text(pos_star[1][0], pos_star[1][1]-0.1, "Single Point of Failure", color=PINK, ha='center')
star_img = save_fig(fig, "star_topology.png")

# 4. Spider-Web Topology
fig, ax = setup_plot(figsize=(6, 6))
nodes = ("M", "S1", "S2", "S3", "S4", "S5", "S6")
edges = [
    ("M", "S1"), ("M", "S2"), ("M", "S3"),
    ("S1", "S2"), ("S2", "S3"),
    ("S1", "S4"), ("S2", "S5"), ("S3", "S6"),
    ("S4", "S5"), ("S5", "S6"),
    ("S2", "S6"), ("S3", "S5")
]
G_spider = nx.Graph()
G_spider.add_nodes_from(nodes)
G_spider.add_edges_from(edges)
pos_spider = nx.spring_layout(G_spider, seed=42)
nx.draw_networkx_nodes(G_spider, pos_spider, node_color=BG_COLOR, edgecolors=CYAN, node_size=1000, linewidths=2, ax=ax)
nx.draw_networkx_edges(G_spider, pos_spider, edge_color=CYAN, style="dashed", width=2, ax=ax)
nx.draw_networkx_labels(G_spider, pos_spider, font_color="white", font_weight="bold", ax=ax)
spider_img = save_fig(fig, "spider_topology.png")

# 5. Architecture Diagram
def draw_flowchart(steps, filename, horizontal=False):
    if horizontal:
        fig, ax = setup_plot(figsize=(10, 2))
        for i, step in enumerate(steps):
            ax.add_patch(patches.Rectangle((i*3, 0), 2.5, 1, facecolor='none', edgecolor=CYAN, linewidth=2))
            ax.text(i*3 + 1.25, 0.5, step, color='white', ha='center', va='center', fontsize=10, fontweight='bold', wrap=True)
            if i < len(steps)-1:
                ax.arrow(i*3 + 2.5, 0.5, 0.3, 0, head_width=0.2, head_length=0.1, fc=CYAN, ec=CYAN)
        ax.set_xlim(-0.5, len(steps)*3)
        ax.set_ylim(-1, 2)
    else:
        fig, ax = setup_plot(figsize=(6, len(steps)*1.5))
        for i, step in enumerate(steps):
            y = len(steps) - i - 1
            ax.add_patch(patches.Rectangle((0, y*2), 5, 1.2, facecolor='none', edgecolor=CYAN, linewidth=2))
            ax.text(2.5, y*2 + 0.6, step, color='white', ha='center', va='center', fontsize=14, fontweight='bold')
            if i < len(steps)-1:
                ax.arrow(2.5, y*2, 0, -0.6, head_width=0.2, head_length=0.2, fc=CYAN, ec=CYAN)
        ax.set_xlim(-1, 6)
        ax.set_ylim(-1, len(steps)*2)
    return save_fig(fig, filename)

arch_img = draw_flowchart(["Network State Monitoring", "Adaptive Parameter Engine", "Smart Route Selection", "Dynamic Pheromone Update", "Packet Delivery"], "architecture.png")
algo_img = draw_flowchart(["Initialize Graph", "Monitor Network", "Compute Adaptive Parameters", "Generate Candidate Paths", "Calculate Path Score", "Select Optimal Route", "Update Pheromone", "Deliver Packet"], "algorithm.png", horizontal=False)

# 6. Math Formulas
def render_math(formula, filename, fontsize=24):
    fig, ax = setup_plot(figsize=(8, 2))
    ax.text(0.5, 0.5, f"${formula}$", color='white', fontsize=fontsize, ha='center', va='center')
    return save_fig(fig, filename)

math_alpha = render_math(r"\alpha_t = \alpha_0 + k_1 \times FailureRate", "math_alpha.png")
math_beta = render_math(r"\beta_t = \beta_0 + k_2 \times AvgCongestion", "math_beta.png")
math_rho = render_math(r"\rho_t = \rho_0 + k_3 \times NetworkInstability", "math_rho.png")
math_instab = render_math(r"Instability = 1 - \frac{AvgReliability}{AvgCongestion}", "math_instab.png")
math_score = render_math(r"Score(P) = \prod_{(u,v)\in P}\tau(u,v)^\alpha \times \left(\frac{r(u,v)}{w(u,v)\times g(u,v)}\right)^\beta", "math_score.png", fontsize=20)
math_update = render_math(r"\tau = (1-\rho)\tau + \frac{R}{L_P}", "math_update.png")
math_pdr = render_math(r"PDR = \frac{PacketsReceived}{PacketsSent}", "math_pdr.png")
math_lat = render_math(r"Latency = \frac{\sum L_i}{m}", "math_lat.png")
math_red = render_math(r"Redundancy = \frac{\sum H_i}{m}", "math_red.png")

# 7. Comparison Bar Charts
def render_bar_chart(title, values, filename):
    fig, ax = plt.subplots(figsize=(6, 4), facecolor=BG_COLOR)
    ax.set_facecolor(BG_COLOR)
    labels = ["Baseline", "Classical ACO", "AS-ACO"]
    colors = ["gray", "white", CYAN]
    bars = ax.bar(labels, values, color=colors)
    ax.set_title(title, color='white', fontsize=16, pad=20)
    ax.tick_params(colors='white')
    for spine in ax.spines.values():
        spine.set_edgecolor('white')
    return save_fig(fig, filename)

chart_pdr = render_bar_chart("Packet Delivery Ratio (PDR)", [0.65, 0.88, 0.98], "chart_pdr.png")
chart_lat = render_bar_chart("Latency (ms)", [45.2, 32.1, 21.5], "chart_lat.png")
chart_red = render_bar_chart("Redundancy (Hops)", [3.8, 2.5, 2.1], "chart_red.png")

# --- BUILD PRESENTATION ---
prs = Presentation()
# Set slide dimensions to widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper to set solid background
def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 15, 26)

def add_title(slide, text):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0, 255, 255) # Cyan
    p.font.bold = True

def add_bullets(slide, items, left=0.5, top=1.8, width=6, height=5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(14)

def add_image(slide, img_path, left, top, width=None, height=None):
    if width:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top))

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Adaptive Smart ACO Based Robust Communication\nfor Vehicle Platooning using Spider-Web Topology"
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(0, 255, 255)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(2))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "By: [Student Name(s)] | Roll No: [Roll No(s)]\nDept. of Electronics and Telecommunication Engineering\n[College Name]\nUnder Guidance of: [Guide Name]"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.alignment = PP_ALIGN.CENTER

# Slide 2: Intro
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Introduction to Vehicle Platooning")
add_bullets(slide, [
    "Platooning connects multiple vehicles in a coordinated convoy.",
    "Driven by a leader, followed tightly by slaves.",
    "Benefits:",
    "  • Reduced fuel consumption",
    "  • Better traffic flow",
    "  • Automated driving coordination"
], width=12)
add_image(slide, platoon_img, 1.5, 4.5, width=10)

# Slide 3: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Problem Statement")
add_bullets(slide, [
    "Communication failure leads to catastrophic consequences:",
    "  • Unsafe spacing",
    "  • Delayed braking",
    "  • Coordination loss",
    "  • High collision risk"
], width=12)
add_image(slide, failure_img, 1.5, 4.5, width=10)

# Slide 4: Existing Model
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Existing Communication Model (Star Topology)")
add_bullets(slide, [
    "Master node directly controls all slaves.",
    "Problem: Single Point of Failure.",
    "If one link breaks, communication is permanently lost.",
    "Lack of path redundancy."
], width=6)
add_image(slide, star_img, 7, 1.5, height=5.5)

# Slide 5: Proposed Model
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Proposed Spider-Web Topology")
add_bullets(slide, [
    "Mesh-like communication links.",
    "Redundant multi-hop architecture.",
    "If a direct link fails, packets route through neighbors.",
    "Ensures structural robustness."
], width=6)
add_image(slide, spider_img, 7, 1.5, height=5.5)

# Slide 6: Classical ACO Limitations
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Why Classical ACO is Insufficient")
add_bullets(slide, [
    "Fixed Alpha, Beta, and Rho parameters.",
    "Cannot adapt dynamically to network congestion.",
    "Slow response during cascading link failures.",
    "Unnecessary latency and premature convergence.",
    "Suboptimal paths in high-stress scenarios."
], width=12)

# Slide 7: AS-ACO Proposed Upgrade
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Proposed Upgrade: Adaptive Smart ACO")
add_bullets(slide, [
    "Introduces real-time dynamic parameter adjustment."
], top=1.5, height=1)
add_image(slide, arch_img, 1, 3.5, width=11)

# Slide 8: Working Principle
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "AS-ACO Working Principle")
add_bullets(slide, [
    "System continuously monitors:",
    "  • Congestion",
    "  • Reliability",
    "  • Link failures",
    "Dynamically adjusts parameters instantly."
], width=5)
add_image(slide, algo_img, 7, 1.5, height=5.5)

# Slide 9: Adaptive Parameter Logic
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Adaptive Parameter Logic")
add_bullets(slide, ["Alpha = Pheromone trust", "Beta = Heuristic influence", "Rho = Pheromone evaporation"])
add_image(slide, math_alpha, 5.5, 2, width=7)
add_image(slide, math_beta, 5.5, 3.5, width=7)
add_image(slide, math_rho, 5.5, 5, width=7)

# Slide 10: Network Monitoring Layer
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Network Monitoring Layer")
add_bullets(slide, [
    "Continuously calculates real-time network health.",
    "Computes Instability Metric dynamically.",
    "Triggers adaptive routing adjustments."
])
add_image(slide, math_instab, 3.5, 4, width=6)

# Slide 11: Route Selection
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Route Selection Algorithm")
add_image(slide, algo_img, 4, 1.5, height=5.5)

# Slide 12: Path Scoring Formula
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Mathematical Path Scoring Model")
add_bullets(slide, [
    "Evaluates candidate paths based on:",
    "  • Pheromone density (τ)",
    "  • Link reliability (r)",
    "  • Latency weight (w)",
    "  • Congestion factor (g)"
])
add_image(slide, math_score, 1, 4, width=11)

# Slide 13: Dynamic Pheromone Update
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Dynamic Pheromone Update")
add_bullets(slide, [
    "Evaporates old pheromones using dynamic Rho.",
    "Deposits fresh pheromones inversely proportional to path length (L_P)."
])
add_image(slide, math_update, 3.5, 4, width=6)

# Slide 14: Performance Metrics
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Evaluation Metrics")
add_image(slide, math_pdr, 1, 2, width=5)
add_image(slide, math_lat, 1, 4, width=5)
add_image(slide, math_red, 1, 6, width=5)
add_bullets(slide, [
    "PDR: High reliability indicator.",
    "Latency: Fast decision making.",
    "Redundancy: Optimization check."
], left=7, width=6)

# Slide 15: Simulation Environment
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Simulation Environment Stack")
add_bullets(slide, [
    "Python 3.10+",
    "NetworkX (Topology & Graphs)",
    "Matplotlib (Real-time Analytics)",
    "Tkinter (Interactive Dashboard)",
    "",
    "Modular Architecture:",
    "  • network.py",
    "  • aco.py",
    "  • simulation.py",
    "  • ui.py"
], width=12)

# Slide 16: Performance Comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Performance Comparison (Table)")
# Create a table
x, y, cx, cy = Inches(1.5), Inches(2.5), Inches(10), Inches(3)
shape = slide.shapes.add_table(4, 4, x, y, cx, cy)
table = shape.table
headers = ["Metric", "Baseline", "Classical ACO", "AS-ACO"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 100, 100)
data = [
    ["Packet Delivery Ratio", "65%", "88%", "98%"],
    ["Average Latency", "45.2 ms", "32.1 ms", "21.5 ms"],
    ["Average Redundancy", "3.8 hops", "2.5 hops", "2.1 hops"]
]
for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 30, 40)

# Slide 17: Graph Comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Performance Benchmarks")
add_image(slide, chart_pdr, 0.5, 2, width=4)
add_image(slide, chart_lat, 4.5, 2, width=4)
add_image(slide, chart_red, 8.5, 2, width=4)

# Slide 18: Advantages
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Advantages of Proposed System")
add_bullets(slide, [
    "✔ Superior Fault Tolerance & Reliability",
    "✔ Real-time Dynamic Routing Intelligence",
    "✔ Zero Single-Point-of-Failure",
    "✔ Instant Response to Congestion",
    "✔ Lower Latency & High PDR"
], width=12)

# Slide 19: Future Scope
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Future Scope")
add_bullets(slide, [
    "• CARLA / SUMO Simulator Integration",
    "• Implementation of Reinforcement Learning (RL)",
    "• 5G V2X Communication Modeling",
    "• Real-world Autonomous Vehicle Integration"
], width=12)

# Slide 20: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Conclusion")
add_bullets(slide, [
    "Adaptive Smart ACO successfully solves the critical limitations of Classical ACO.",
    "By dynamically adjusting its routing behavior based on real-time network conditions, the system ensures maximum reliability.",
    "This paves the way for highly intelligent and robust communication in next-generation autonomous vehicle platooning."
], width=12)

# Slide 21: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You\nAny Questions?"
p.font.size = Pt(64)
p.font.color.rgb = RGBColor(0, 255, 255)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

prs.save("Final_Project_Presentation.pptx")
print("Presentation generated successfully!")
